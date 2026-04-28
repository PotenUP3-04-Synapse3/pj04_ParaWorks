from __future__ import annotations

import io
from typing import Any

import structlog

from backend.parsers.base import BaseParser, ParsedDocument

log = structlog.get_logger(__name__)


class PdfParser(BaseParser):
    def can_parse(self, mime_type: str) -> bool:
        return mime_type == 'application/pdf'

    def parse(self, content: bytes | str, source_url: str | None = None) -> ParsedDocument:
        from pypdf import PdfReader  # type: ignore
        if isinstance(content, str):
            content = content.encode('utf-8')

        reader = PdfReader(io.BytesIO(content))
        pages_text: list[str] = []
        paragraphs: list[dict[str, Any]] = []
        para_idx = 0

        for page_num, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or '').strip()
            pages_text.append(text)
            for line in text.splitlines():
                line = line.strip()
                if line:
                    paragraphs.append({
                        'text': line,
                        'page_number': page_num,
                        'paragraph_index': para_idx,
                        'heading': None,
                    })
                    para_idx += 1

        return ParsedDocument(
            text='\n'.join(pages_text),
            pages=pages_text,
            paragraphs=paragraphs,
            metadata={'parser': 'pypdf', 'num_pages': len(reader.pages)},
        )


class OfficeParser(BaseParser):
    """Word (.docx), Excel (.xlsx), PowerPoint (.pptx) 파서."""

    SUPPORTED = {
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    }

    def can_parse(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED

    def parse(self, content: bytes | str, source_url: str | None = None) -> ParsedDocument:
        if isinstance(content, str):
            content = content.encode('utf-8')
        # 실제 MIME 타입 추론은 호출측에서 전달
        # 여기서는 magic bytes로 분기
        buf = content[:8]
        if buf[:2] == b'PK':  # ZIP-based Office formats
            return self._parse_ooxml(content)
        return ParsedDocument(text='[지원되지 않는 Office 파일 형식]', metadata={'parser': 'office_unsupported'})

    def _parse_ooxml(self, content: bytes) -> ParsedDocument:
        # MIME 타입 추론을 위해 python-docx/openpyxl/python-pptx를 순서대로 시도
        paragraphs: list[dict[str, Any]] = []
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []

        # Word
        try:
            from docx import Document  # type: ignore
            doc = Document(io.BytesIO(content))
            para_idx = 0
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    paragraphs.append({'text': t, 'page_number': None, 'paragraph_index': para_idx, 'heading': para.style.name if 'Heading' in para.style.name else None})
                    text_parts.append(t)
                    para_idx += 1
            for table in doc.tables:
                tbl = [[cell.text for cell in row.cells] for row in table.rows]
                tables.append(tbl)
                text_parts.append('\n'.join('\t'.join(row) for row in tbl))
            if text_parts:
                return ParsedDocument(text='\n'.join(text_parts), paragraphs=paragraphs, tables=tables, metadata={'parser': 'python-docx'})
        except Exception:
            pass

        # Excel
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            rows: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_str = '\t'.join(str(c) if c is not None else '' for c in row)
                    rows.append(row_str)
            full = '\n'.join(rows)
            if full.strip():
                return ParsedDocument(text=full, metadata={'parser': 'openpyxl'})
        except Exception:
            pass

        # PowerPoint
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(io.BytesIO(content))
            slides_text: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_texts = [shape.text_frame.text for shape in slide.shapes if hasattr(shape, 'text_frame')]
                slides_text.append(f'[슬라이드 {i}]\n' + '\n'.join(slide_texts))
            full = '\n\n'.join(slides_text)
            if full.strip():
                return ParsedDocument(text=full, metadata={'parser': 'python-pptx'})
        except Exception:
            pass

        return ParsedDocument(text='[Office 파일 파싱 실패]', metadata={'parser': 'office_failed'})


class TextParser(BaseParser):
    SUPPORTED = {'text/plain', 'text/markdown', 'text/html', 'text/csv'}

    def can_parse(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED

    def parse(self, content: bytes | str, source_url: str | None = None) -> ParsedDocument:
        if isinstance(content, bytes):
            content = content.decode('utf-8', errors='replace')

        # HTML → plain text
        if 'html' in (source_url or '') or content.strip().startswith('<'):
            from bs4 import BeautifulSoup  # type: ignore
            content = BeautifulSoup(content, 'lxml').get_text(separator='\n')

        lines = [l.strip() for l in content.splitlines() if l.strip()]
        paragraphs = [
            {'text': line, 'page_number': None, 'paragraph_index': i, 'heading': None}
            for i, line in enumerate(lines)
        ]
        return ParsedDocument(
            text='\n'.join(lines),
            paragraphs=paragraphs,
            metadata={'parser': 'text'},
        )

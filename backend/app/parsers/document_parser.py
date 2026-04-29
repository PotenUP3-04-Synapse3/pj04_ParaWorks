"""File parsers for all supported document types."""
from __future__ import annotations

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text_from_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(f'[표] {row_text}')
    return '\n'.join(parts)


def extract_text_from_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            parts.append(text)
    return '\n'.join(parts)


def extract_text_from_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            parts.append(f'[슬라이드 {i}]\n' + '\n'.join(slide_texts))
    return '\n\n'.join(parts)


def extract_text_from_xlsx(data: bytes) -> str:
    import pandas as pd
    parts = []
    xls = pd.ExcelFile(io.BytesIO(data))
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        parts.append(f'[시트: {sheet_name}]\n{df.to_csv(index=False)}')
    return '\n\n'.join(parts)


def extract_text_from_image(data: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(img, lang='kor+eng')
    except Exception as exc:
        logger.warning('OCR failed: %s', exc)
        return ''


def extract_text(data: bytes, mime_type: str, file_name: str = '') -> str:
    """Dispatch to the appropriate parser based on mime type or extension."""
    ext = Path(file_name).suffix.lower() if file_name else ''

    if mime_type in ('application/vnd.openxmlformats-officedocument.wordprocessingml.document',) or ext == '.docx':
        return extract_text_from_docx(data)

    if mime_type == 'application/pdf' or ext == '.pdf':
        return extract_text_from_pdf(data)

    if mime_type in ('application/vnd.openxmlformats-officedocument.presentationml.presentation',) or ext == '.pptx':
        return extract_text_from_pptx(data)

    if mime_type in ('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',) or ext in ('.xlsx', '.xls'):
        return extract_text_from_xlsx(data)

    if ext in ('.hwp', '.hwpx'):
        from app.parsers.hwp_parser import extract_text_from_hwp_bytes
        return extract_text_from_hwp_bytes(data, suffix=ext)

    if mime_type.startswith('image/') or ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'):
        return extract_text_from_image(data)

    if mime_type in ('text/plain', 'text/markdown') or ext in ('.txt', '.md'):
        return data.decode('utf-8', errors='replace')

    if mime_type == 'text/csv' or ext == '.csv':
        return data.decode('utf-8', errors='replace')

    logger.warning('No parser available for mime_type=%s ext=%s', mime_type, ext)
    return ''
